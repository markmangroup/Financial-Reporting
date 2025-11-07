#!/usr/bin/env python3
"""
PowerPoint Content Extractor & Analyzer
Extracts slides, text, and creates analysis for audit comparison
"""

import os
import sys
from pathlib import Path
import json

def install_requirements():
    """Install required packages if not available"""
    try:
        import pip
        packages = [
            'python-pptx',
            'Pillow',
            'matplotlib',
            'pandas'
        ]

        for package in packages:
            try:
                __import__(package.replace('-', '_'))
                print(f"✅ {package} already installed")
            except ImportError:
                print(f"📦 Installing {package}...")
                pip.main(['install', package])
    except Exception as e:
        print(f"⚠️  Package installation issue: {e}")
        print("Please run: pip install python-pptx Pillow matplotlib pandas")

def extract_powerpoint_content(pptx_path):
    """Extract all content from PowerPoint file"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        import io

        # Create output directory
        output_dir = Path("powerpoint-analysis")
        output_dir.mkdir(exist_ok=True)

        # Load presentation
        prs = Presentation(pptx_path)
        print(f"🎯 Found {len(prs.slides)} slides to analyze")

        analysis_data = {
            "total_slides": len(prs.slides),
            "slides": [],
            "all_text": "",
            "charts_found": [],
            "numbers_found": []
        }

        # Extract content from each slide
        for i, slide in enumerate(prs.slides, 1):
            print(f"📊 Processing slide {i}...")

            slide_data = {
                "slide_number": i,
                "title": "",
                "text_content": [],
                "shapes": [],
                "charts": [],
                "numbers": []
            }

            # Process all shapes on slide
            for shape in slide.shapes:
                shape_info = {"type": str(shape.shape_type)}

                # Extract text content
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_data["text_content"].append(text)
                    analysis_data["all_text"] += f" {text}"

                    # Set slide title (usually first text shape)
                    if not slide_data["title"] and len(text) < 100:
                        slide_data["title"] = text

                    # Extract numbers
                    import re
                    numbers = re.findall(r'[\$]?[\d,]+\.?\d*[%]?', text)
                    if numbers:
                        slide_data["numbers"].extend(numbers)
                        analysis_data["numbers_found"].extend(numbers)

                # Identify chart types
                if "CHART" in str(shape.shape_type):
                    chart_info = f"Chart on slide {i}"
                    slide_data["charts"].append(chart_info)
                    analysis_data["charts_found"].append(chart_info)
                    shape_info["chart"] = True

                slide_data["shapes"].append(shape_info)

            analysis_data["slides"].append(slide_data)

        # Save analysis to JSON
        with open(output_dir / "analysis.json", "w") as f:
            json.dump(analysis_data, f, indent=2)

        # Create summary report
        create_analysis_report(analysis_data, output_dir)

        print(f"✅ Analysis complete! Check the '{output_dir}' folder")
        return analysis_data

    except ImportError as e:
        print(f"❌ Missing required library: {e}")
        print("Installing required packages...")
        install_requirements()
        return extract_powerpoint_content(pptx_path)
    except Exception as e:
        print(f"❌ Error processing PowerPoint: {e}")
        return None

def create_analysis_report(data, output_dir):
    """Create detailed analysis report"""
    report_path = output_dir / "audit_report.md"

    with open(report_path, "w") as f:
        f.write("# PowerPoint Analysis Report\n\n")
        f.write(f"**Total Slides:** {data['total_slides']}\n")
        f.write(f"**Charts Found:** {len(data['charts_found'])}\n")
        f.write(f"**Numbers Extracted:** {len(data['numbers_found'])}\n\n")

        f.write("## Slide-by-Slide Analysis\n\n")
        for slide in data['slides']:
            f.write(f"### Slide {slide['slide_number']}\n")
            f.write(f"**Title:** {slide['title']}\n\n")

            if slide['text_content']:
                f.write("**Content:**\n")
                for text in slide['text_content']:
                    f.write(f"- {text}\n")
                f.write("\n")

            if slide['charts']:
                f.write("**Charts:**\n")
                for chart in slide['charts']:
                    f.write(f"- {chart}\n")
                f.write("\n")

            if slide['numbers']:
                f.write("**Key Numbers:**\n")
                for num in slide['numbers']:
                    f.write(f"- {num}\n")
                f.write("\n")

            f.write("---\n\n")

        f.write("## All Extracted Numbers\n\n")
        unique_numbers = list(set(data['numbers_found']))
        for num in sorted(unique_numbers):
            f.write(f"- {num}\n")

        f.write("\n## Charts Identified\n\n")
        for chart in data['charts_found']:
            f.write(f"- {chart}\n")

def reconcile_with_our_data():
    """Create reconciliation template with our golden record"""
    reconciliation = {
        "our_numbers": {
            "total_revenue": 181320.01,
            "total_expenses": 312800.27,
            "net_loss": -131480.26,
            "current_cash": 48941.01,
            "consultant_spend": 125081.72,
            "credit_card_expenses": 152374.31,
            "auto_loan_expenses": 34034.24,
            "bank_fees": 1310.00
        },
        "powerpoint_numbers": {},
        "variances": {},
        "notes": "Numbers to be filled in after PowerPoint analysis"
    }

    with open("powerpoint-analysis/reconciliation.json", "w") as f:
        json.dump(reconciliation, f, indent=2)

    print("📋 Created reconciliation template")

def main():
    print("🚀 PowerPoint Analysis & Audit Tool")
    print("=" * 50)

    # Find PowerPoint file
    pptx_path = "data/Strategic & Operational Overview.pptx"
    if not os.path.exists(pptx_path):
        print(f"❌ PowerPoint file not found: {pptx_path}")
        return

    print(f"📁 Found PowerPoint file: {pptx_path}")
    print(f"📏 File size: {os.path.getsize(pptx_path) / 1024 / 1024:.1f} MB")

    # Extract and analyze content
    analysis = extract_powerpoint_content(pptx_path)

    if analysis:
        # Create reconciliation framework
        reconcile_with_our_data()

        print("\n🎯 NEXT STEPS:")
        print("1. Check 'powerpoint-analysis/audit_report.md' for detailed analysis")
        print("2. Review 'powerpoint-analysis/analysis.json' for structured data")
        print("3. Use 'powerpoint-analysis/reconciliation.json' for number comparison")

        # Quick summary
        print(f"\n📊 QUICK SUMMARY:")
        print(f"   • {analysis['total_slides']} slides processed")
        print(f"   • {len(analysis['charts_found'])} charts identified")
        print(f"   • {len(set(analysis['numbers_found']))} unique numbers extracted")

        return True

    return False

if __name__ == "__main__":
    success = main()
    if success:
        print("\n✅ Analysis complete! Check the output files for detailed insights.")
    else:
        print("\n❌ Analysis failed. Check the error messages above.")